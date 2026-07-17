import os
import fitz  # PyMuPDF
from pptx import Presentation
from langchain.schema import Document

def parse_pdf(file_path: str, base_metadata: dict) -> list[Document]:
    documents = []
    try:
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    meta = base_metadata.copy()
                    meta["page_number"] = page_num + 1
                    documents.append(Document(page_content=text, metadata=meta))
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return documents

def parse_ppt(file_path: str, base_metadata: dict) -> list[Document]:
    documents = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text = text.strip()
            if text:
                meta = base_metadata.copy()
                meta["slide_number"] = slide_num + 1
                documents.append(Document(page_content=text, metadata=meta))
    except Exception as e:
        print(f"Error reading PPT: {e}")
    return documents

def extract_documents(file_path: str, file_type: str, metadata: dict) -> list[Document]:
    if file_type.lower() == 'pdf':
        return parse_pdf(file_path, metadata)
    elif file_type.lower() in ['ppt', 'pptx']:
        return parse_ppt(file_path, metadata)
    else:
        # Fallback to reading as text
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read().strip()
                if text:
                    return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            print(f"Error reading text file: {e}")
            
    return []
